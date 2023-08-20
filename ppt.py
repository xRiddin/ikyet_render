import json, os
from models.generate import generate_response as ge
# from IKYET.ikyet_render.models.sdxl import gen as img
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
prs = Presentation()


async def content(topics):
    con = await ge("""you are expert at following the given prompt. I would like to prepare a speech in pure markdown format, 
                with language and topic consistent with the outline provided. 
                You do not need to add any additional replies or explanations.
                Here is the global outline:
                {topics}
                you will generate content for each page accordingly. 
                Each title must corresponds to one page of content, with a maximum of 100 words per page that is size of a small essay.
                you must only reply in python list of strings for each topic you generate, strictly follow this format:
                [{title:(title 1 generated), content:(content 1 generated)}, {title:(title 2 generated), content:(content 2 generated)}]
                """, f" please return a minimum of 10 pages of content about {topics}", )

    return con


async def slides(topics, direc):
        while True:
            try:
                content0 = await content(topics)
                print("this is content" + content0)
                content1 = json.loads(content0)
                for slide_data in content1:
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    background = slide.background
                    fill = background.fill
                    fill.solid()
                    fill.fore_color.rgb = RGBColor(0, 0, 0)
                    background_image_path = "static/i1.jpg"

                # Calculate the dimensions for the background image
                    left = top = Inches(0)
                    width = prs.slide_width
                    height = prs.slide_height
                    background_image = slide.shapes.add_picture(background_image_path, left, top, width, height)
                # Send the image to the back
                    slide.shapes._spTree.remove(background_image._element)
                    slide.shapes._spTree.insert(2, background_image._element)

                    title_shape = slide.shapes.title
                    title_shape.text = slide_data["title"]

                    text_box = slide.shapes.placeholders[1]
                    text_box.text = slide_data["content"]
                os.makedirs(direc)
                prs.save(f"{direc}/generated.pptx")
                return f'{direc}/generated.pptx'
            except Exception as e:
                print(e)
        

# implement image creation

